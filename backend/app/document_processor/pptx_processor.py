import logging
import pptx

logger = logging.getLogger(__name__)


class PPTXProcessor:
    """Extracts text from PowerPoint presentations (.pptx) slide by slide."""
    
    @staticmethod
    def extract(file_path: str) -> list[dict]:
        try:
            prs = pptx.Presentation(file_path)
        except Exception as e:
            raise ValueError(f"Failed to read PowerPoint document (.pptx): {str(e)}")
        
        slides_data = []
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_text = []
            slide_title = f"Slide {slide_idx}"
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        if shape == slide.shapes.title and text:
                            slide_title = text
                        slide_text.append(text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            slide_text.append(row_text)
            
            full_text = "\n".join(slide_text).strip()
            if full_text:
                slides_data.append({
                    'page_number': slide_idx,
                    'section': slide_title,
                    'slide': slide_idx,
                    'text': full_text,
                    'file_type': 'pptx'
                })
        
        if not slides_data:
            raise ValueError("No text content found in PowerPoint presentation.")
            
        logger.info(f"Extracted {len(slides_data)} slides from PPTX {file_path}")
        return slides_data
